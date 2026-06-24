#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Exporta o MESMO contrato slides.json em .pptx NATIVO e editavel (python-pptx).

Frente 2 do estudio de apresentacoes: fecha a lacuna "saida editavel" vs. Gamma.
Le `slides.json` + `marca/tokens.json` (a mesma fonte que `gerar.py` consome para
o PDF) e emite um .pptx onde cada slide vira FORMAS/PLACEHOLDERS NATIVOS do
PowerPoint - caixas de texto, tabelas, autoshapes e CHART nativo - e NAO imagens.
O dono abre no PowerPoint/Google Slides e edita o texto, move as caixas, ajusta
as cores. Aditivo: NAO altera o caminho PDF (`gerar.py`/templates).

Mapeamento tipo -> forma nativa (os 10 tipos do contrato):
  capa       -> slide de titulo (titulo/subtitulo/autor/data)
  topicos    -> titulo + bullets reais (paragrafos de lista, nao imagem)
  destaque   -> numero grande ambar + legenda
  texto      -> titulo + corpo (html degradado para texto, sem renderizar HTML)
  comparacao -> duas caixas lado a lado; coluna "novo" (direita) em ambar
  timeline   -> marcos em sequencia (linha + bolinhas + caixas ano/texto)
  checklist  -> itens numerados (1., 2., 3. ...)
  processo   -> etapas em caixas com setas (autoshapes) entre elas
  grafico    -> grafico de BARRAS NATIVO (XL_CHART_TYPE.COLUMN_CLUSTERED) com a
                regua de honestidade: eixo de valor com MINIMO 0, maximo = campo
                `maximo` quando presente, rotulos de valor visiveis, e COR COMO
                EXCECAO (so a barra `destaque` em ambar; o resto em navy).
                Ancora: MemoryCode/conhecimento/sinteses/
                o-grafico-convence-antes-do-numero.md (base no zero, cor como
                excecao - o grafico convence antes do numero).
  kpis       -> 2 a 4 caixas numero+legenda lado a lado.
Campos comuns: intro (sob o titulo) e rodape_nota (fonte/ressalva no pe).

A marca vem de tokens.json (navy #1F3A5F + amber #E08A3C), nada cravado.

Uso:
  python3 gerar-pptx.py [--dados exemplos/deck.json] [--saida saida/deck.pptx]

Dependencia: python-pptx (requirements.txt). Sem ela, aborta com erro claro.
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

import tema as _tema

RAIZ = Path(__file__).resolve().parent
DIR_MARCA = RAIZ / "marca"
DIR_SAIDA = RAIZ / "saida"
TOKENS_JSON = DIR_MARCA / "tokens.json"

# Geometria 16:9 em polegadas (python-pptx mede em Emu/Inches/Pt).
LARGURA_SLIDE_IN = 13.333
ALTURA_SLIDE_IN = 7.5
MARGEM_IN = 0.9  # margem horizontal das caixas de conteudo
TOPO_TITULO_IN = 0.55
TOPO_CONTEUDO_IN = 1.9  # onde o corpo comeca quando ha titulo (+intro)


# ---------------------------------------------------------------------------
# Marca: cores e tipografia a partir de tokens.json (fonte unica, igual ao PDF)
# ---------------------------------------------------------------------------
def _rgb(hexcor: str):
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(hexcor.lstrip("#").upper())


def _primeira_familia(pilha_css: str) -> str:
    """'Georgia, "Times New Roman", serif' -> 'Georgia'. python-pptx aplica UMA
    familia por run; o Office cai no fallback se a fonte faltar no sistema."""
    primeira = pilha_css.split(",")[0].strip()
    return primeira.strip("'").strip('"')


class Marca:
    """Espelha os tokens de cor/tipografia em objetos do python-pptx. So le o que
    o deck precisa; a fonte da verdade continua em tokens.json (o mesmo do PDF)."""

    def __init__(self, tokens: dict) -> None:
        cor = tokens.get("cor", {})
        tipo = tokens.get("tipografia", {})
        marca = tokens.get("marca", {})

        self.primaria = _rgb(cor.get("primaria", "#1F3A5F"))
        self.primaria_clara = _rgb(cor.get("primaria_clara", "#345B8C"))
        # Fundo da CAPA: por padrao a primaria (capa navy de hoje). Um tema pode
        # sobrescrever com `capa_fundo` quando a primaria deixa de servir de fundo
        # (no escuro a primaria virou azul claro -> branco sobre ela seria ilegivel,
        # entao capa_fundo da o fundo escuro). Sem a chave, capa = primaria (claro
        # identico). A cor do TITULO/marca na capa segue clara (branco/acento).
        self.capa_fundo = _rgb(cor.get("capa_fundo", cor.get("primaria", "#1F3A5F")))
        # Subtitulo/autor da capa: por padrao neutro_300 (cinza claro sobre o navy
        # de hoje = identico). Num tema com capa dedicada, segue 'capa_subtitulo'
        # para nao herdar um neutro_300 que ficou ESCURO (ilegivel na capa escura).
        self.capa_subtitulo = _rgb(
            cor.get("capa_subtitulo", cor.get("neutro_300", "#C9CED6"))
        )
        self.secundaria = _rgb(cor.get("secundaria", "#3E7CB1"))
        self.acento = _rgb(cor.get("acento", "#E08A3C"))
        self.texto = _rgb(cor.get("texto", "#1A1D21"))
        self.texto_suave = _rgb(cor.get("texto_suave", "#6B7280"))
        self.fundo = _rgb(cor.get("fundo", "#FFFFFF"))
        self.fundo_suave = _rgb(cor.get("fundo_suave", "#EEF1F5"))
        self.neutro_300 = _rgb(cor.get("neutro_300", "#C9CED6"))
        self.branco = _rgb("#FFFFFF")

        self.fonte_display = _primeira_familia(tipo.get("fonte_display", "Georgia"))
        self.fonte_corpo = _primeira_familia(
            tipo.get("fonte_corpo", "Helvetica Neue, Arial, sans-serif")
        )
        self.nome = marca.get("nome", "")
        self.assinatura_rodape = marca.get("assinatura_rodape", "")


# ---------------------------------------------------------------------------
# Helpers de texto/forma (todos produzem objetos NATIVOS, editaveis no Office)
# ---------------------------------------------------------------------------
def _caixa_texto(slide, esq, topo, larg, alt):
    from pptx.util import Inches

    caixa = slide.shapes.add_textbox(
        Inches(esq), Inches(topo), Inches(larg), Inches(alt)
    )
    tf = caixa.text_frame
    tf.word_wrap = True
    return caixa, tf


def _set_run(run, texto, *, tamanho, cor, fonte, negrito=False):
    from pptx.util import Pt

    run.text = texto
    run.font.size = Pt(tamanho)
    run.font.bold = negrito
    run.font.name = fonte
    run.font.color.rgb = cor


def _paragrafo(tf, texto, *, tamanho, cor, fonte, negrito=False,
               novo=True, espaco_antes=None, alinhar=None):
    """Adiciona paragrafo formatado. O 1o paragrafo do tf ja existe vazio, entao
    `novo=False` reusa em vez de criar uma linha em branco no topo."""
    from pptx.util import Pt

    p = tf.add_paragraph() if novo else tf.paragraphs[0]
    if espaco_antes is not None:
        p.space_before = Pt(espaco_antes)
    if alinhar is not None:
        p.alignment = alinhar
    run = p.add_run()
    _set_run(run, texto, tamanho=tamanho, cor=cor, fonte=fonte, negrito=negrito)
    return p


def _retangulo(slide, esq, topo, larg, alt, cor_fundo, *,
               cor_linha=None, arredondado=False):
    """Autoshape retangular (preenchida) nativa - editavel/movível no Office."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    forma = MSO_SHAPE.ROUNDED_RECTANGLE if arredondado else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        forma, Inches(esq), Inches(topo), Inches(larg), Inches(alt)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = cor_fundo
    if cor_linha is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = cor_linha
        shape.line.width = _pt(1)
    shape.shadow.inherit = False  # sem sombra: dado-tinta alto (Tufte)
    return shape


def _pt(valor):
    from pptx.util import Pt

    return Pt(valor)


def _strip_html(html: str) -> list[str]:
    """Degrada HTML do campo `texto.html` para paragrafos de texto puro, SEM
    renderizar HTML: separa blocos por </p>/<br>, remove tags, desescapa entidades
    comuns. O contrato e 'sem renderizar HTML' - aqui ele vira texto editavel."""
    import html as _html

    bruto = re.sub(r"(?i)</p\s*>|<br\s*/?>", "\n", html)
    sem_tags = re.sub(r"<[^>]+>", "", bruto)
    texto = _html.unescape(sem_tags)
    return [linha.strip() for linha in texto.splitlines() if linha.strip()]


# ---------------------------------------------------------------------------
# Blocos compartilhados: titulo, intro, rodape de marca, rodape_nota
# ---------------------------------------------------------------------------
def _add_titulo(slide, m: Marca, titulo: str) -> None:
    _, tf = _caixa_texto(slide, MARGEM_IN, TOPO_TITULO_IN,
                         LARGURA_SLIDE_IN - 2 * MARGEM_IN, 1.0)
    _paragrafo(tf, titulo, tamanho=30, cor=m.primaria, fonte=m.fonte_display,
               negrito=True, novo=False)
    # Filete ambar sob o titulo (acento da marca, dado-tinta minimo).
    _retangulo(slide, MARGEM_IN, TOPO_TITULO_IN + 0.95, 1.4, 0.05, m.acento)


def _add_intro(slide, m: Marca, intro: str) -> float:
    """Paragrafo de contexto sob o titulo. Devolve o topo (in) onde o corpo deve
    comecar, empurrado para baixo quando ha intro."""
    if not intro:
        return TOPO_CONTEUDO_IN
    _, tf = _caixa_texto(slide, MARGEM_IN, TOPO_TITULO_IN + 1.05,
                         LARGURA_SLIDE_IN - 2 * MARGEM_IN, 0.7)
    _paragrafo(tf, intro, tamanho=14, cor=m.texto_suave, fonte=m.fonte_corpo,
               novo=False)
    return TOPO_CONTEUDO_IN + 0.55


def _add_rodape_nota(slide, m: Marca, nota: str) -> None:
    if not nota:
        return
    _, tf = _caixa_texto(slide, MARGEM_IN, ALTURA_SLIDE_IN - 0.85,
                         LARGURA_SLIDE_IN - 2 * MARGEM_IN, 0.4)
    _paragrafo(tf, nota, tamanho=10, cor=m.texto_suave, fonte=m.fonte_corpo,
               novo=False)


def _add_rodape_marca(slide, m: Marca) -> None:
    """Assinatura discreta no canto - identidade consistente em todo slide."""
    from pptx.enum.text import PP_ALIGN

    if not m.assinatura_rodape:
        return
    _, tf = _caixa_texto(slide, MARGEM_IN, ALTURA_SLIDE_IN - 0.45,
                         LARGURA_SLIDE_IN - 2 * MARGEM_IN, 0.3)
    _paragrafo(tf, m.assinatura_rodape, tamanho=8, cor=m.neutro_300,
               fonte=m.fonte_corpo, novo=False, alinhar=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Renderizadores por tipo de slide
# ---------------------------------------------------------------------------
def render_capa(slide, m: Marca, deck: dict, dados: dict) -> None:
    # Faixa de fundo cheia (capa = bloco de marca). No claro = navy primaria; no
    # escuro = capa_fundo (fundo escuro), para o titulo branco ficar legivel.
    fundo = _retangulo(slide, 0, 0, LARGURA_SLIDE_IN, ALTURA_SLIDE_IN, m.capa_fundo)
    fundo.line.fill.background()

    if m.nome:
        _, tf = _caixa_texto(slide, MARGEM_IN, 1.0,
                             LARGURA_SLIDE_IN - 2 * MARGEM_IN, 0.5)
        _paragrafo(tf, m.nome.upper(), tamanho=14, cor=m.acento,
                   fonte=m.fonte_corpo, negrito=True, novo=False)

    titulo = dados.get("titulo") or deck.get("titulo", "")
    _, tf = _caixa_texto(slide, MARGEM_IN, 2.4,
                         LARGURA_SLIDE_IN - 2 * MARGEM_IN, 2.2)
    _paragrafo(tf, titulo, tamanho=44, cor=m.branco, fonte=m.fonte_display,
               negrito=True, novo=False)

    subtitulo = dados.get("subtitulo") or deck.get("subtitulo", "")
    if subtitulo:
        _paragrafo(tf, subtitulo, tamanho=20, cor=m.capa_subtitulo,
                   fonte=m.fonte_corpo, espaco_antes=14)

    # Autor / data no pe.
    autor = dados.get("autor") or deck.get("autor", "")
    data = dados.get("data") or deck.get("data", "")
    linha = " · ".join(p for p in (autor, data) if p)
    if linha:
        _, tf = _caixa_texto(slide, MARGEM_IN, ALTURA_SLIDE_IN - 1.0,
                             LARGURA_SLIDE_IN - 2 * MARGEM_IN, 0.5)
        _paragrafo(tf, linha, tamanho=13, cor=m.capa_subtitulo, fonte=m.fonte_corpo,
                   novo=False)


def render_topicos(slide, m: Marca, dados: dict) -> None:
    _add_titulo(slide, m, dados.get("titulo", ""))
    topo = _add_intro(slide, m, dados.get("intro", ""))
    _, tf = _caixa_texto(slide, MARGEM_IN, topo,
                         LARGURA_SLIDE_IN - 2 * MARGEM_IN,
                         ALTURA_SLIDE_IN - topo - 1.0)
    pontos = dados.get("pontos", [])
    for i, ponto in enumerate(pontos):
        p = _paragrafo(tf, "", tamanho=20, cor=m.texto, fonte=m.fonte_corpo,
                       novo=(i > 0), espaco_antes=(0 if i == 0 else 12))
        # Marcador ambar (run separado) + texto navy: cor como acento pontual.
        _set_run(p.runs[0], "▪  ", tamanho=20, cor=m.acento,
                 fonte=m.fonte_corpo)
        run_txt = p.add_run()
        _set_run(run_txt, str(ponto), tamanho=20, cor=m.texto, fonte=m.fonte_corpo)


def render_destaque(slide, m: Marca, dados: dict) -> None:
    from pptx.enum.text import PP_ALIGN

    # Numero grande ambar centralizado + legenda navy.
    _, tf = _caixa_texto(slide, MARGEM_IN, 2.0,
                         LARGURA_SLIDE_IN - 2 * MARGEM_IN, 2.0)
    _paragrafo(tf, str(dados.get("numero", "")), tamanho=88, cor=m.acento,
               fonte=m.fonte_display, negrito=True, novo=False,
               alinhar=PP_ALIGN.CENTER)

    legenda = dados.get("legenda", "")
    if legenda:
        _, tf2 = _caixa_texto(slide, 2.0, 4.5, LARGURA_SLIDE_IN - 4.0, 1.5)
        _paragrafo(tf2, legenda, tamanho=22, cor=m.primaria, fonte=m.fonte_corpo,
                   novo=False, alinhar=PP_ALIGN.CENTER)


def render_texto(slide, m: Marca, dados: dict) -> None:
    _add_titulo(slide, m, dados.get("titulo", ""))
    topo = _add_intro(slide, m, dados.get("intro", ""))
    paragrafos = _strip_html(dados.get("html", ""))
    _, tf = _caixa_texto(slide, MARGEM_IN, topo,
                         LARGURA_SLIDE_IN - 2 * MARGEM_IN,
                         ALTURA_SLIDE_IN - topo - 1.0)
    for i, par in enumerate(paragrafos):
        _paragrafo(tf, par, tamanho=18, cor=m.texto, fonte=m.fonte_corpo,
                   novo=(i > 0), espaco_antes=(0 if i == 0 else 10))


def _coluna_comparacao(slide, m: Marca, esq, topo, larg, alt, coluna, destacar):
    """Uma coluna da comparacao: caixa com sub-titulo + itens. A coluna 'novo'
    (direita) ganha fundo ambar suave e barra ambar - cor como mensagem."""
    cor_borda = m.acento if destacar else m.neutro_300
    cor_titulo = m.acento if destacar else m.primaria
    _retangulo(slide, esq, topo, larg, alt, m.fundo_suave,
               cor_linha=cor_borda, arredondado=True)

    _, tf = _caixa_texto(slide, esq + 0.25, topo + 0.2, larg - 0.5, alt - 0.4)
    _paragrafo(tf, coluna.get("titulo", ""), tamanho=18, cor=cor_titulo,
               fonte=m.fonte_display, negrito=True, novo=False)
    for item in coluna.get("itens", []):
        p = _paragrafo(tf, "", tamanho=15, cor=m.texto, fonte=m.fonte_corpo,
                       espaco_antes=10)
        _set_run(p.runs[0], "•  ", tamanho=15,
                 cor=(m.acento if destacar else m.secundaria), fonte=m.fonte_corpo)
        run_txt = p.add_run()
        _set_run(run_txt, str(item), tamanho=15, cor=m.texto, fonte=m.fonte_corpo)


def render_comparacao(slide, m: Marca, dados: dict) -> None:
    _add_titulo(slide, m, dados.get("titulo", ""))
    topo = _add_intro(slide, m, dados.get("intro", ""))
    alt = ALTURA_SLIDE_IN - topo - 1.0
    larg_total = LARGURA_SLIDE_IN - 2 * MARGEM_IN
    larg_col = (larg_total - 0.5) / 2  # 0.5in de gap entre colunas
    _coluna_comparacao(slide, m, MARGEM_IN, topo, larg_col, alt,
                       dados.get("esquerda", {}), destacar=False)
    _coluna_comparacao(slide, m, MARGEM_IN + larg_col + 0.5, topo, larg_col, alt,
                       dados.get("direita", {}), destacar=True)


def render_timeline(slide, m: Marca, dados: dict) -> None:
    from pptx.enum.text import PP_ALIGN

    _add_titulo(slide, m, dados.get("titulo", ""))
    topo = _add_intro(slide, m, dados.get("intro", ""))
    marcos = dados.get("marcos", [])
    if not marcos:
        return
    larg_total = LARGURA_SLIDE_IN - 2 * MARGEM_IN
    n = len(marcos)
    passo = larg_total / n
    y_linha = topo + 1.1

    # Linha horizontal navy (autoshape fina) ligando os marcos.
    _retangulo(slide, MARGEM_IN, y_linha, larg_total, 0.04, m.primaria_clara)

    for i, marco in enumerate(marcos):
        cx = MARGEM_IN + passo * i + passo / 2
        # Bolinha ambar no eixo.
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        d = 0.28
        bola = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(y_linha - d / 2 + 0.02),
            Inches(d), Inches(d)
        )
        bola.fill.solid()
        bola.fill.fore_color.rgb = m.acento
        bola.line.color.rgb = m.fundo
        bola.line.width = _pt(1.5)
        bola.shadow.inherit = False

        # Ano acima da linha.
        _, tf_ano = _caixa_texto(slide, cx - passo / 2 + 0.1, y_linha - 0.85,
                                 passo - 0.2, 0.6)
        _paragrafo(tf_ano, str(marco.get("ano", "")), tamanho=16, cor=m.primaria,
                   fonte=m.fonte_display, negrito=True, novo=False,
                   alinhar=PP_ALIGN.CENTER)
        # Texto abaixo da linha.
        _, tf_txt = _caixa_texto(slide, cx - passo / 2 + 0.1, y_linha + 0.35,
                                 passo - 0.2, 2.0)
        _paragrafo(tf_txt, str(marco.get("texto", "")), tamanho=13, cor=m.texto,
                   fonte=m.fonte_corpo, novo=False, alinhar=PP_ALIGN.CENTER)


def render_checklist(slide, m: Marca, dados: dict) -> None:
    _add_titulo(slide, m, dados.get("titulo", ""))
    topo = _add_intro(slide, m, dados.get("intro", ""))
    itens = dados.get("itens", [])
    _, tf = _caixa_texto(slide, MARGEM_IN, topo,
                         LARGURA_SLIDE_IN - 2 * MARGEM_IN,
                         ALTURA_SLIDE_IN - topo - 1.0)
    for i, item in enumerate(itens):
        p = _paragrafo(tf, "", tamanho=19, cor=m.texto, fonte=m.fonte_corpo,
                       novo=(i > 0), espaco_antes=(0 if i == 0 else 14))
        _set_run(p.runs[0], f"{i + 1}.  ", tamanho=19, cor=m.acento,
                 fonte=m.fonte_display, negrito=True)
        run_txt = p.add_run()
        _set_run(run_txt, str(item), tamanho=19, cor=m.texto, fonte=m.fonte_corpo)


def render_processo(slide, m: Marca, dados: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches

    _add_titulo(slide, m, dados.get("titulo", ""))
    topo = _add_intro(slide, m, dados.get("intro", ""))
    etapas = dados.get("etapas", [])
    if not etapas:
        return
    n = len(etapas)
    larg_total = LARGURA_SLIDE_IN - 2 * MARGEM_IN
    gap = 0.45  # espaco para a seta entre caixas
    larg_caixa = (larg_total - gap * (n - 1)) / n
    y = topo + 0.4
    alt_caixa = 2.4

    for i, etapa in enumerate(etapas):
        x = MARGEM_IN + (larg_caixa + gap) * i
        caixa = _retangulo(slide, x, y, larg_caixa, alt_caixa, m.fundo_suave,
                           cor_linha=m.neutro_300, arredondado=True)
        tf = caixa.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        _paragrafo(tf, str(etapa.get("rotulo", "")), tamanho=17, cor=m.primaria,
                   fonte=m.fonte_display, negrito=True, novo=False,
                   alinhar=PP_ALIGN.CENTER)
        _paragrafo(tf, str(etapa.get("texto", "")), tamanho=13, cor=m.texto,
                   fonte=m.fonte_corpo, espaco_antes=8, alinhar=PP_ALIGN.CENTER)

        # Seta ambar entre esta caixa e a proxima.
        if i < n - 1:
            seta = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + larg_caixa + 0.05), Inches(y + alt_caixa / 2 - 0.15),
                Inches(gap - 0.1), Inches(0.3)
            )
            seta.fill.solid()
            seta.fill.fore_color.rgb = m.acento
            seta.line.fill.background()
            seta.shadow.inherit = False


def _como_numero(valor):
    """Converte para float quando der (aceita o "5"/"3.8" que o JSON pode trazer
    e que o render ja coage com float()); devolve None quando nao e numero. Evita
    o `valor > maximo` cru, que estoura (TypeError) ao misturar str e num e
    compara string lexicograficamente ("10" > "5" e False) - falso negativo."""
    try:
        return float(valor)
    except (TypeError, ValueError):
        return None


def _avisar_estouro_grafico(dados: dict) -> None:
    """Quando o slide declara `maximo` (teto de dominio) e alguma barra o excede,
    o eixo fixo em `maximo` faz a barra encostar/passar do topo: e honesto com a
    ESCALA mas esconde que o dado estourou o teto. Avisa em stderr (nao quebra o
    render) cada barra que estoura. Sem `maximo` o eixo e auto, sem estouro."""
    maximo = _como_numero(dados.get("maximo"))
    if maximo is None:
        return
    titulo = dados.get("titulo", "")
    for barra in dados.get("barras", []):
        valor = _como_numero(barra.get("valor"))
        if valor is not None and valor > maximo:
            print(
                f"[aviso] grafico '{titulo}': valor {barra.get('valor')} de "
                f"'{barra.get('rotulo', '')}' excede o maximo declarado ({maximo:g})"
                " - barra limitada a 100%, escala pode enganar",
                file=sys.stderr,
            )


def render_grafico(slide, m: Marca, dados: dict) -> None:
    """Grafico de barras NATIVO (editavel: dados em planilha embutida no .pptx).

    Regua de honestidade (ancora: o-grafico-convence-antes-do-numero.md):
      - eixo de valor com MINIMO 0 (nunca trunca a base) -> minimum_scale = 0;
      - maximo = campo `maximo` quando presente (escala de dominio); senao auto;
      - rotulos de valor visiveis em cada barra (data labels);
      - COR COMO EXCECAO: so a barra `destaque` em ambar; o resto em navy.
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
    from pptx.util import Inches, Pt

    _avisar_estouro_grafico(dados)
    _add_titulo(slide, m, dados.get("titulo", ""))
    topo = _add_intro(slide, m, dados.get("intro", ""))
    barras = dados.get("barras", [])
    if not barras:
        return

    rotulos = [str(b.get("rotulo", "")) for b in barras]
    # coage com seguranca: "3.8" -> 3.8, valor invalido ("abc") -> 0 (nao crasha
    # como o float() cru fazia). Alinhado ao web/PDF e ao _avisar_estouro_grafico.
    valores = [_como_numero(b.get("valor")) or 0 for b in barras]

    chart_data = CategoryChartData()
    chart_data.categories = rotulos
    unidade = dados.get("unidade", "")
    chart_data.add_series(unidade or "Valor", valores)

    larg = LARGURA_SLIDE_IN - 2 * MARGEM_IN
    alt = ALTURA_SLIDE_IN - topo - 1.2
    grafico_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(MARGEM_IN), Inches(topo), Inches(larg), Inches(alt),
        chart_data,
    )
    chart = grafico_frame.chart
    chart.has_legend = False  # serie unica: legenda seria tinta redundante (Tufte)
    chart.has_title = False

    # --- Eixo de valor: HONESTIDADE. Minimo SEMPRE 0; maximo = `maximo` ou auto.
    eixo_valor = chart.value_axis
    eixo_valor.minimum_scale = 0.0
    _maximo = _como_numero(dados.get("maximo"))
    if _maximo is not None:
        eixo_valor.maximum_scale = _maximo
    eixo_valor.has_major_gridlines = True
    eixo_valor.major_gridlines.format.line.color.rgb = m.neutro_300
    eixo_valor.major_gridlines.format.line.width = Pt(0.5)
    eixo_valor.tick_labels.font.size = Pt(11)
    eixo_valor.tick_labels.font.color.rgb = m.texto_suave

    eixo_cat = chart.category_axis
    eixo_cat.tick_labels.font.size = Pt(12)
    eixo_cat.tick_labels.font.color.rgb = m.texto

    plot = chart.plots[0]
    plot.gap_width = 80
    serie = plot.series[0]

    # Rotulos de valor visiveis em cada barra (le-se o numero sem cruzar o eixo).
    serie.has_data_labels = True
    dls = serie.data_labels
    dls.show_value = True  # sem isto o XML sai com showVal=0 e o numero nao aparece
    dls.number_format = "0.0" if any(v != int(v) for v in valores) else "0"
    dls.number_format_is_linked = False
    dls.position = XL_LABEL_POSITION.OUTSIDE_END
    dls.font.size = Pt(11)
    dls.font.bold = True
    dls.font.color.rgb = m.primaria

    # COR COMO EXCECAO: navy por padrao; ambar so na(s) barra(s) `destaque`.
    serie.format.fill.solid()
    serie.format.fill.fore_color.rgb = m.primaria
    for i, b in enumerate(barras):
        if b.get("destaque"):
            ponto = serie.points[i]
            ponto.format.fill.solid()
            ponto.format.fill.fore_color.rgb = m.acento


def render_kpis(slide, m: Marca, dados: dict) -> None:
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    _add_titulo(slide, m, dados.get("titulo", ""))
    topo = _add_intro(slide, m, dados.get("intro", ""))
    itens = dados.get("itens", [])[:4]  # contrato: 2 a 4 KPIs
    if not itens:
        return
    n = len(itens)
    larg_total = LARGURA_SLIDE_IN - 2 * MARGEM_IN
    gap = 0.4
    larg_caixa = (larg_total - gap * (n - 1)) / n
    y = topo + 0.5
    alt_caixa = 2.6

    for i, item in enumerate(itens):
        x = MARGEM_IN + (larg_caixa + gap) * i
        caixa = _retangulo(slide, x, y, larg_caixa, alt_caixa, m.fundo_suave,
                           cor_linha=m.neutro_300, arredondado=True)
        tf = caixa.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _paragrafo(tf, str(item.get("numero", "")), tamanho=48, cor=m.acento,
                   fonte=m.fonte_display, negrito=True, novo=False,
                   alinhar=PP_ALIGN.CENTER)
        _paragrafo(tf, str(item.get("legenda", "")), tamanho=14, cor=m.primaria,
                   fonte=m.fonte_corpo, espaco_antes=10, alinhar=PP_ALIGN.CENTER)


# Mapa tipo -> renderizador. Capa e tratada a parte (precisa do bloco deck).
RENDERIZADORES = {
    "topicos": render_topicos,
    "destaque": render_destaque,
    "texto": render_texto,
    "comparacao": render_comparacao,
    "timeline": render_timeline,
    "checklist": render_checklist,
    "processo": render_processo,
    "grafico": render_grafico,
    "kpis": render_kpis,
}


# ---------------------------------------------------------------------------
# Montagem do deck
# ---------------------------------------------------------------------------
def montar_pptx(dados_path: Path, saida: Path, tema: str | None = None) -> int:
    """Le slides.json + tokens.json e grava o .pptx. Devolve o numero de slides
    escritos (a capa conta como 1). `tema` faz o restyle por cor: so a paleta da
    Marca muda; as formas/posicoes/tipografia de estrutura ficam iguais."""
    from pptx import Presentation
    from pptx.util import Inches

    bruto = json.loads(dados_path.read_text(encoding="utf-8"))
    deck = bruto.get("deck", {})
    slides = bruto.get("slides", [])
    tokens = json.loads(TOKENS_JSON.read_text(encoding="utf-8"))
    m = Marca(_tema.aplicar_tema(tokens, tema))  # so 'cor' muda com o tema

    pres = Presentation()
    pres.slide_width = Inches(LARGURA_SLIDE_IN)
    pres.slide_height = Inches(ALTURA_SLIDE_IN)
    layout_branco = pres.slide_layouts[6]  # layout em branco: montamos as formas

    n_escritos = 0

    # Capa: sempre o primeiro slide, a partir do bloco `deck` (ou do 1o slide
    # tipo 'capa', se o contrato trouxer um explicito).
    capa_explicita = next((s for s in slides if s.get("tipo") == "capa"), None)
    slide_capa = pres.slides.add_slide(layout_branco)
    render_capa(slide_capa, m, deck, capa_explicita or {})
    n_escritos += 1

    nao_suportados: list[str] = []
    for s in slides:
        tipo = s.get("tipo")
        if tipo == "capa":
            continue  # ja virou o slide de abertura
        render = RENDERIZADORES.get(tipo)
        if render is None:
            nao_suportados.append(str(tipo))
            continue
        slide = pres.slides.add_slide(layout_branco)
        render(slide, m, s)
        _add_rodape_nota(slide, m, s.get("rodape_nota", ""))
        _add_rodape_marca(slide, m)
        n_escritos += 1

    saida.parent.mkdir(parents=True, exist_ok=True)
    pres.save(str(saida))

    print(f"[gerar-pptx] PPTX gerado: {saida} ({n_escritos} slides).")
    if nao_suportados:
        print(f"[gerar-pptx] AVISO: tipos ignorados (sem renderizador): "
              f"{', '.join(sorted(set(nao_suportados)))}", file=sys.stderr)
    return n_escritos


def _resolver_saida(caminho: Path) -> Path:
    """Confina a escrita a esta pasta (mesma regra do gerar.py): bloqueia '..'
    e caminho absoluto fora da arvore, sem impedir --saida relativo legitimo."""
    destino = (caminho if caminho.is_absolute() else RAIZ / caminho).resolve()
    try:
        destino.relative_to(RAIZ)
    except ValueError:
        print(f"[gerar-pptx] ERRO: --saida deve ficar dentro de {RAIZ} "
              f"(recebido: {destino}).", file=sys.stderr)
        sys.exit(2)
    return destino


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Exporta slides.json em .pptx nativo editavel (python-pptx)."
    )
    parser.add_argument("--dados", type=Path,
                        help="JSON do deck (default: exemplos/deck.json)")
    parser.add_argument("--saida", type=Path,
                        help="caminho do .pptx (default: saida/deck.pptx)")
    parser.add_argument(
        "--tema", default=None,
        help="tema de cor (restyle): ausente ou 'claro' = paleta atual; "
             "'escuro' = fundo escuro acessivel. So muda cor, nao o layout.",
    )
    args = parser.parse_args(argv)

    # Valida o tema cedo (erro claro com a lista, nunca silencioso).
    tokens_check = json.loads(TOKENS_JSON.read_text(encoding="utf-8"))
    try:
        _tema.resolver_cor(tokens_check, args.tema)
    except _tema.TemaInexistente as e:
        print(f"[gerar-pptx] ERRO: {e}", file=sys.stderr)
        return 2

    try:
        import pptx  # noqa: F401
    except ImportError:
        print("[gerar-pptx] ERRO: dependencia 'python-pptx' ausente.",
              file=sys.stderr)
        print(f"[gerar-pptx] Instale: pip install -r {RAIZ / 'requirements.txt'}",
              file=sys.stderr)
        return 2

    dados = args.dados or RAIZ / "exemplos" / "deck.json"
    saida = _resolver_saida(args.saida or DIR_SAIDA / "deck.pptx")
    montar_pptx(dados, saida, args.tema)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
